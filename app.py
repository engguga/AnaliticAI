import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches
import io
import requests

# --- Configuração da página ---
st.set_page_config(page_title="Analitic", layout="wide")
st.title("📊 Analitic – Relatórios Inteligentes com IA Hugging Face")

# --- Token via sidebar ---
HF_TOKEN = st.sidebar.text_input("🔑 Hugging Face API Token", type="password")

# --- Upload do Excel ---
uploaded_file = st.file_uploader("📎 Envie sua planilha Excel", type=[".xlsx"])

# --- Função IA Hugging Face ---
def gerar_insight_hf(prompt, token):
    API_URL = "https://api-inference.huggingface.co/models/google/flan-t5-small"
    headers = {"Authorization": f"Bearer {token}"}
    payload = {"inputs": prompt}
    try:
        response = requests.post(API_URL, headers=headers, json=payload, timeout=30)
        if response.status_code == 200:
            return response.json()[0]["generated_text"]
        else:
            return f"Erro ({response.status_code}): {response.text}"
    except Exception as e:
        return f"Erro na chamada Hugging Face: {e}"

# --- Lógica principal ---
if uploaded_file:
    df = pd.read_excel(uploaded_file)
    st.success("Arquivo carregado com sucesso!")

    st.subheader("📋 Prévia dos Dados")
    df_display = df.copy()
    if 'Data_Admissao' in df_display.columns:
        df_display['Data_Admissao'] = df_display['Data_Admissao'].dt.strftime('%Y-%m-%d')
    st.dataframe(df_display.head())

    st.subheader("📈 Estatísticas Descritivas")
    stats_desc = df.describe().round(2)
    st.dataframe(stats_desc)

    # Gráficos
    if 'Idade' in df.columns:
        st.subheader("📊 Distribuição de Idades")
        fig_idade, ax = plt.subplots()
        sns.histplot(df['Idade'], kde=True, ax=ax, color='skyblue')
        ax.set_title("Distribuição de Idades")
        st.pyplot(fig_idade)
        plt.close(fig_idade)

    if {'Salario', 'Departamento'}.issubset(df.columns):
        st.subheader("📊 Salário por Departamento")
        fig_salario, ax = plt.subplots()
        sns.boxplot(x='Departamento', y='Salario', data=df, ax=ax, palette='viridis')
        ax.set_title("Salário por Departamento")
        plt.xticks(rotation=45)
        st.pyplot(fig_salario)
        plt.close(fig_salario)

    # Insights manuais
    st.subheader("💡 Insights Estatísticos")
    insights = []
    if 'Salario' in df.columns:
        media_salario = df['Salario'].mean().round(2)
        insights.append(f"A média salarial geral é de: ${media_salario}")
    if 'Departamento' in df.columns and 'Salario' in df.columns:
        salario_medio_por_dept = df.groupby('Departamento')['Salario'].mean().sort_values(ascending=False).round(2)
        insights.append(f"O departamento com maior média salarial é: {salario_medio_por_dept.index[0]} (${salario_medio_por_dept.iloc[0]})")
        insights.append(f"O departamento com menor média salarial é: {salario_medio_por_dept.index[-1]} (${salario_medio_por_dept.iloc[-1]})")
    for i in insights:
        st.markdown(f"- {i}")

    # Insight via IA Hugging Face
    st.subheader("🤖 Insight com IA Hugging Face")
    insight_ia = ""
    if HF_TOKEN:
        prompt = f"Explique os dados a seguir de forma profissional para um relatório executivo:\n\n{stats_desc.to_string()}\n\nConclusão:"
        with st.spinner("Chamando IA da Hugging Face..."):
            insight_ia = gerar_insight_hf(prompt, HF_TOKEN)
            st.success("Insight gerado com sucesso!")
            st.markdown(insight_ia)
    else:
        st.warning("Insira seu Hugging Face Token na barra lateral para gerar o insight com IA.")

    # --- PowerPoint ---
    if st.button("📤 Gerar Apresentação PowerPoint"):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Relatório Analítico"
        slide.placeholders[1].text = "Gerado com Analitic + IA Hugging Face"

        # Estatísticas
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Estatísticas Descritivas"
        tf = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4)).text_frame
        tf.text = stats_desc.to_string()

        # Gráfico de Idades
        if 'Idade' in df.columns:
            buf = io.BytesIO()
            fig_idade, ax = plt.subplots()
            sns.histplot(df['Idade'], kde=True, ax=ax, color='skyblue')
            ax.set_title("Distribuição de Idades")
            plt.savefig(buf)
            plt.close(fig_idade)
            buf.seek(0)
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "Distribuição de Idades"
            slide.shapes.add_picture(buf, Inches(1), Inches(2), width=Inches(8))

        # Gráfico de Salário
        if {'Salario', 'Departamento'}.issubset(df.columns):
            buf = io.BytesIO()
            fig_salario, ax = plt.subplots()
            sns.boxplot(x='Departamento', y='Salario', data=df, ax=ax, palette='viridis')
            ax.set_title("Salário por Departamento")
            plt.xticks(rotation=45)
            plt.savefig(buf, bbox_inches='tight')
            plt.close(fig_salario)
            buf.seek(0)
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "Salário por Departamento"
            slide.shapes.add_picture(buf, Inches(1), Inches(2), width=Inches(8))

        # Slide com IA
        if insight_ia:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Insight com IA"
            slide.placeholders[1].text = insight_ia

        # Download
        buf_pptx = io.BytesIO()
        prs.save(buf_pptx)
        st.download_button(
            label="⬇️ Baixar Apresentação PowerPoint",
            data=buf_pptx.getvalue(),
            file_name="relatorio_analitic_ia.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )